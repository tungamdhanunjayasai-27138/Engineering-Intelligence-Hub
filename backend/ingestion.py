from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from backend.chunking import chunk_text
from backend.embeddings import create_embedding

DOCUMENTS_DIR = Path("data/documents")
DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)


def read_file(file):
    ext = file.suffix.lower()

    if ext == ".pdf":
        reader = PdfReader(file)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if ext == ".docx":
        doc = Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext in [".txt", ".md", ".html", ".htm"]:
        return file.read_text(encoding="utf-8")

    if ext == ".csv":
        return file.read_text(encoding="utf-8")

    if ext == ".xlsx":
        wb = load_workbook(file, read_only=True, data_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                text.append(" | ".join(str(x) for x in row if x is not None))
        return "\n".join(text)

    if ext == ".pptx":
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    raise ValueError(f"Unsupported file type: {ext}")


def load_documents():
    documents = []

    for file in DOCUMENTS_DIR.iterdir():
        if file.is_file():
            try:
                documents.append({
                    "name": file.name,
                    "text": read_file(file)
                })
            except ValueError:
                continue

    return documents


def process_documents():
    results = []

    for document in load_documents():
        chunks = chunk_text(document["text"])

        for i, chunk in enumerate(chunks):
            results.append({
                "source": document["name"],
                "chunk_id": i,
                "text": chunk,
                "embedding": create_embedding(chunk)
            })

    return results